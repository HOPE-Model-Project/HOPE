from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
PRESENTATION_DIR = Path(__file__).resolve().parent
OUTPUT_PPTX = PRESENTATION_DIR / "GERMANY_PCM_CASE_BUILD_DECK.pptx"
OUTPUT_MD = PRESENTATION_DIR / "GERMANY_PCM_CASE_BUILD_DECK.md"

NODAL_CASE = ROOT / "ModelCases" / "GERMANY_PCM_nodal_case" / "Data_GERMANY_PCM_nodal"
ZONAL_CASE = ROOT / "ModelCases" / "GERMANY_PCM_zonal4_case" / "Data_GERMANY_PCM_zonal4"
DEMO_SYSTEM_COST = ROOT / "ModelCases" / "GERMANY_PCM_nodal_jan_2day_rescaled_case" / "output" / "system_cost.csv"
GERMANY_TOOLS = ROOT / "tools" / "germany_pcm_case_related"
GERMANY_DASHBOARD_GEOJSON = ROOT / "tools" / "hope_dashboard" / "data" / "germany_tso_zones.geojson"

COLOR_BG = RGBColor(246, 248, 251)
COLOR_PANEL = RGBColor(255, 255, 255)
COLOR_NAVY = RGBColor(18, 42, 76)
COLOR_BLUE = RGBColor(47, 93, 169)
COLOR_TEAL = RGBColor(54, 143, 132)
COLOR_GOLD = RGBColor(219, 165, 32)
COLOR_CORAL = RGBColor(225, 107, 74)
COLOR_PURPLE = RGBColor(145, 96, 192)
COLOR_SLATE = RGBColor(91, 105, 125)
COLOR_LINE = RGBColor(214, 223, 233)
COLOR_WHITE = RGBColor(255, 255, 255)

ZONE_COLORS = [COLOR_BLUE, COLOR_CORAL, COLOR_TEAL, COLOR_PURPLE]


def _case_metrics() -> dict[str, object]:
    nodal_buses = len(pd.read_csv(NODAL_CASE / "busdata.csv"))
    nodal_lines = len(pd.read_csv(NODAL_CASE / "linedata.csv"))
    nodal_gens = len(pd.read_csv(NODAL_CASE / "gendata.csv"))
    chronology_rows = len(pd.read_csv(NODAL_CASE / "load_timeseries_nodal.csv"))
    zonal_zones = len(pd.read_csv(ZONAL_CASE / "zonedata.csv"))
    zonal_lines = len(pd.read_csv(ZONAL_CASE / "linedata.csv"))
    zonal_gens = len(pd.read_csv(ZONAL_CASE / "gendata.csv"))
    demo_cost = pd.read_csv(DEMO_SYSTEM_COST)
    total_demo_cost = float(demo_cost["Total_cost ($)"].sum())
    total_demo_lol = float(demo_cost["LoL_plt ($)"].sum())
    zone_costs = [
        (str(row["Zone"]), float(row["Opr_cost ($)"]) / 1e6)
        for _, row in demo_cost.iterrows()
    ]
    return {
        "nodal_buses": nodal_buses,
        "nodal_lines": nodal_lines,
        "nodal_gens": nodal_gens,
        "chronology_rows": chronology_rows,
        "zonal_zones": zonal_zones,
        "zonal_lines": zonal_lines,
        "zonal_gens": zonal_gens,
        "demo_cost_m": total_demo_cost / 1e6,
        "demo_lol": total_demo_lol,
        "zone_costs": zone_costs,
    }


def _set_background(slide, color: RGBColor = COLOR_BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left: float, top: float, width: float, height: float, text: str = ""):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if text:
        box.text_frame.text = text
    return box


def _style_runs(paragraph, size: int, color: RGBColor, bold: bool = False) -> None:
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = _textbox(slide, 0.6, 0.38, 11.8, 0.55, title)
    p = title_box.text_frame.paragraphs[0]
    _style_runs(p, 26, COLOR_NAVY, bold=True)
    if subtitle:
        sub_box = _textbox(slide, 0.62, 0.96, 11.4, 0.4, subtitle)
        p = sub_box.text_frame.paragraphs[0]
        _style_runs(p, 13, COLOR_SLATE)


def _add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = COLOR_PANEL, line: RGBColor = COLOR_LINE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def _card_title(shape, text: str, size: int = 12, color: RGBColor = COLOR_SLATE) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    _style_runs(p, size, color, bold=True)


def _add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], font_size: int = 18) -> None:
    box = _textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        _style_runs(p, font_size, COLOR_NAVY)


def _add_metric_card(slide, left: float, top: float, width: float, label: str, value: str, accent: RGBColor) -> None:
    card = _add_card(slide, left, top, width, 1.0, fill=COLOR_WHITE)
    card.line.color.rgb = accent
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.14),
        Inches(1.0),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.color.rgb = accent

    label_box = _textbox(slide, left + 0.24, top + 0.14, width - 0.34, 0.25, label)
    _style_runs(label_box.text_frame.paragraphs[0], 11, COLOR_SLATE, bold=True)
    value_box = _textbox(slide, left + 0.24, top + 0.42, width - 0.34, 0.34, value)
    _style_runs(value_box.text_frame.paragraphs[0], 24, COLOR_NAVY, bold=True)


def _add_title_slide(prs: Presentation, metrics: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.55))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_NAVY
    band.line.color.rgb = COLOR_NAVY

    title = _textbox(slide, 0.72, 0.42, 8.0, 0.6, "Germany PCM Case Build")
    _style_runs(title.text_frame.paragraphs[0], 28, COLOR_WHITE, bold=True)
    subtitle = _textbox(slide, 0.74, 1.02, 7.4, 0.3, "From raw public data to consistent nodal and zonal HOPE cases")
    _style_runs(subtitle.text_frame.paragraphs[0], 14, RGBColor(222, 231, 242))

    quote = _add_card(slide, 0.78, 1.95, 5.2, 1.15, fill=RGBColor(236, 242, 250))
    quote.text_frame.clear()
    p1 = quote.text_frame.paragraphs[0]
    p1.text = "Core design rule"
    _style_runs(p1, 12, COLOR_BLUE, bold=True)
    p2 = quote.text_frame.add_paragraph()
    p2.text = "Build one nodal master case first, then derive the zonal case mechanically."
    _style_runs(p2, 20, COLOR_NAVY, bold=True)

    _add_metric_card(slide, 0.8, 3.55, 2.45, "Nodal Buses", f"{metrics['nodal_buses']}", COLOR_BLUE)
    _add_metric_card(slide, 3.45, 3.55, 2.45, "Nodal Lines", f"{metrics['nodal_lines']}", COLOR_CORAL)
    _add_metric_card(slide, 6.10, 3.55, 2.45, "Generators", f"{metrics['nodal_gens']}", COLOR_TEAL)
    _add_metric_card(slide, 8.75, 3.55, 2.45, "Demo Cost", f"${metrics['demo_cost_m']:.2f}M", COLOR_GOLD)

    footer = _textbox(slide, 0.8, 5.15, 11.4, 1.35)
    tf = footer.text_frame
    tf.clear()
    for idx, line in enumerate(
        [
            "Network backbone: OSM Europe transmission dataset + PyPSA-Eur reference workflow",
            "Fleet: powerplantmatching with Germany-specific validation layers",
            "Chronology: SMARD national load + generation and four TSO load helper files",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        _style_runs(p, 17, COLOR_NAVY)
        p.space_after = Pt(8)


def _add_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Executive Summary", "How the Germany case was assembled and why the design is comparison-friendly")

    cards = [
        ("Goal", "Create one canonical Germany PCM dataset that supports both nodal and zonal analysis.", COLOR_BLUE),
        ("Method", "Freeze every major mapping once: generator -> bus -> zone -> zonal derivative.", COLOR_TEAL),
        ("Outcome", "A solved 2-day nodal demo case plus a consistent 4-zone zonal comparison case.", COLOR_CORAL),
    ]
    for idx, (title, body, color) in enumerate(cards):
        left = 0.78 + idx * 4.13
        card = _add_card(slide, left, 1.65, 3.55, 3.1)
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(1.65), Inches(3.55), Inches(0.16))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.color.rgb = color
        label = _textbox(slide, left + 0.22, 1.95, 2.9, 0.3, title)
        _style_runs(label.text_frame.paragraphs[0], 15, color, bold=True)
        body_box = _textbox(slide, left + 0.22, 2.35, 3.0, 1.9, body)
        _style_runs(body_box.text_frame.paragraphs[0], 22, COLOR_NAVY, bold=True)


def _add_source_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Source Stack", "Each layer was chosen to keep the zonal case derived from, not detached from, the nodal case")

    cards = [
        ("Network", "OSM Europe transmission dataset\nPyPSA-Eur workflow reference", COLOR_BLUE),
        ("Fleet", "powerplantmatching\nBNetzA validation support", COLOR_CORAL),
        ("Chronology", "SMARD Germany load + generation\nSMARD TSO load helper files", COLOR_TEAL),
        ("Map Geometry", "Germany state boundary GeoJSON\nreconstructed 4-TSO overlay", COLOR_PURPLE),
    ]
    positions = [(0.85, 1.65), (6.8, 1.65), (0.85, 4.0), (6.8, 4.0)]
    for (title, body, color), (left, top) in zip(cards, positions):
        card = _add_card(slide, left, top, 5.55, 1.85)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.22), Inches(top + 0.22), Inches(1.3), Inches(0.38))
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.color.rgb = color
        chip.text_frame.text = title
        chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _style_runs(chip.text_frame.paragraphs[0], 11, COLOR_WHITE, bold=True)
        body_box = _textbox(slide, left + 0.24, top + 0.72, 4.95, 0.9, body)
        for p in body_box.text_frame.paragraphs:
            _style_runs(p, 18, COLOR_NAVY, bold=True)


def _add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Integration Architecture", "The zonal case is produced by aggregation from the nodal master case")

    boxes = [
        ("Raw Sources", 0.8, COLOR_BLUE),
        ("Cleaned Staging", 3.0, COLOR_CORAL),
        ("Frozen Maps", 5.2, COLOR_TEAL),
        ("Nodal Master", 7.4, COLOR_BLUE),
        ("Zonal Derivative", 9.6, COLOR_CORAL),
    ]
    for label, left, color in boxes:
        shape = _add_card(slide, left, 2.0, 1.8, 0.95, fill=color, line=color)
        shape.text_frame.text = label
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _style_runs(shape.text_frame.paragraphs[0], 16, COLOR_WHITE, bold=True)

    for left in [2.65, 4.85, 7.05, 9.25]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(left), Inches(2.26), Inches(0.35), Inches(0.38))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_SLATE
        arrow.line.color.rgb = COLOR_SLATE

    bullets = [
        "Network cleaner produces stable buses, lines, transformers, and links.",
        "Fleet cleaner assigns each generator to one nodal bus and inherits one zone tag.",
        "Chronology builder freezes one 2025 hourly basis and optional zonal load shares.",
        "Zonal seams and capacities are derived from the nodal cross-zone structure.",
    ]
    _add_bullets(slide, 0.9, 3.5, 11.4, 2.8, bullets, font_size=18)


def _add_steps_slide(prs: Presentation, metrics: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Build Steps", "Six steps turned raw public data into the Germany PCM comparison set")

    steps = [
        ("1", "Network backbone", f"Normalize OSM Europe inputs and keep a connected {metrics['nodal_buses']}-bus Germany network."),
        ("2", "Bus-zone mapping", "Freeze one Bus_id -> TSO zone table for 50Hertz, Amprion, TenneT, and TransnetBW."),
        ("3", "Generator mapping", "Clean powerplantmatching and assign each plant to a nodal bus once."),
        ("4", "Chronology", f"Build one 2025 hourly chronology with {metrics['chronology_rows']} rows and reuse it everywhere."),
        ("5", "Case assembly", "Build the canonical nodal case first and aggregate the zonal case from it."),
        ("6", "Dashboard geometry", "Reconstruct a file-backed Germany TSO overlay from state boundaries and frozen bus geography."),
    ]
    top = 1.55
    for num, title, body in steps:
        num_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), Inches(top), Inches(0.55), Inches(0.55))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = COLOR_NAVY
        num_shape.line.color.rgb = COLOR_NAVY
        num_shape.text_frame.text = num
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _style_runs(num_shape.text_frame.paragraphs[0], 18, COLOR_WHITE, bold=True)

        title_box = _textbox(slide, 1.7, top - 0.02, 3.0, 0.28, title)
        _style_runs(title_box.text_frame.paragraphs[0], 18, COLOR_BLUE, bold=True)
        body_box = _textbox(slide, 1.72, top + 0.24, 10.0, 0.42, body)
        _style_runs(body_box.text_frame.paragraphs[0], 16, COLOR_NAVY)
        top += 0.93


def _add_snapshot_slide(prs: Presentation, metrics: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Case Snapshot", "Current master and derivative cases in the repo")

    left_card = _add_card(slide, 0.85, 1.65, 5.7, 4.2)
    right_card = _add_card(slide, 6.8, 1.65, 5.7, 4.2)

    left_title = _textbox(slide, 1.1, 1.92, 2.0, 0.25, "Nodal Master Case")
    _style_runs(left_title.text_frame.paragraphs[0], 18, COLOR_BLUE, bold=True)
    right_title = _textbox(slide, 7.05, 1.92, 2.3, 0.25, "Zonal Derivative Case")
    _style_runs(right_title.text_frame.paragraphs[0], 18, COLOR_CORAL, bold=True)

    left_bullets = [
        f"{metrics['nodal_buses']} buses",
        f"{metrics['nodal_lines']} lines",
        f"{metrics['nodal_gens']} generator entries",
        "8760-hour canonical chronology",
    ]
    right_bullets = [
        f"{metrics['zonal_zones']} zones",
        f"{metrics['zonal_lines']} interzonal interfaces",
        f"{metrics['zonal_gens']} zonal generator entries",
        "Derived from nodal geography and seams",
    ]
    _add_bullets(slide, 1.1, 2.35, 4.8, 2.3, left_bullets, font_size=19)
    _add_bullets(slide, 7.05, 2.35, 4.8, 2.3, right_bullets, font_size=19)

    note = _add_card(slide, 1.1, 4.9, 10.95, 0.62, fill=RGBColor(236, 242, 250))
    note.text_frame.text = "Design choice: the zonal case is not independently calibrated; it is aggregated from the nodal master dataset."
    _style_runs(note.text_frame.paragraphs[0], 15, COLOR_NAVY, bold=True)


def _add_assumptions_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Assumptions And Caveats", "Important choices to explain when comparing Germany nodal and zonal results")

    left = _add_card(slide, 0.85, 1.65, 5.8, 4.6)
    right = _add_card(slide, 6.75, 1.65, 5.75, 4.6)

    lt = _textbox(slide, 1.08, 1.9, 2.6, 0.25, "Frozen Modeling Choices")
    _style_runs(lt.text_frame.paragraphs[0], 18, COLOR_TEAL, bold=True)
    rt = _textbox(slide, 6.98, 1.9, 2.3, 0.25, "Current Caveats")
    _style_runs(rt.text_frame.paragraphs[0], 18, COLOR_CORAL, bold=True)

    left_bullets = [
        "Use one nodal master case as the canonical truth source.",
        "Use four TSO research zones instead of the real DE-LU bidding zone.",
        "Keep representative-day modes off for PCM debug validation windows.",
        "Preserve chronology consistency before adding more model complexity.",
    ]
    right_bullets = [
        "The Germany TSO dashboard overlay is reconstructed, not official GIS.",
        "The 2-day nodal demo case is solve-proven; longer nodal horizons still need staged validation.",
        "Generator validation against MaStR and Kraftwerksliste can still be deepened.",
        "Historical benchmarking should keep expanding as the case matures.",
    ]
    _add_bullets(slide, 1.1, 2.35, 5.1, 3.6, left_bullets, font_size=18)
    _add_bullets(slide, 7.0, 2.35, 5.0, 3.6, right_bullets, font_size=18)


def _add_debug_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Debugging Journey", "The nodal case became usable after fixing topology and electrical scaling")

    milestones = [
        ("Initial failure", "Disconnected components, singular-network behavior, and severe angle-limit pathologies.", COLOR_CORAL),
        ("Core fixes", "Transformer connectivity, largest connected component, and HOPE-scale reactance normalization.", COLOR_TEAL),
        ("Usable result", "A 1-day and then 2-day full-nodal Germany debug run solved cleanly with zero load shedding.", COLOR_BLUE),
    ]
    for idx, (title, body, color) in enumerate(milestones):
        left = 0.9 + idx * 4.1
        card = _add_card(slide, left, 2.0, 3.55, 2.9)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(2.0), Inches(3.55), Inches(0.16))
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.color.rgb = color
        title_box = _textbox(slide, left + 0.22, 2.28, 2.8, 0.26, title)
        _style_runs(title_box.text_frame.paragraphs[0], 17, color, bold=True)
        body_box = _textbox(slide, left + 0.22, 2.7, 2.95, 1.7, body)
        _style_runs(body_box.text_frame.paragraphs[0], 17, COLOR_NAVY)


def _add_results_slide(prs: Presentation, metrics: dict[str, object]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Solved Demo Outcome", "2-day January nodal case used for the current Germany dashboard demo")

    _add_metric_card(slide, 0.85, 1.55, 2.55, "Solve Window", "48 hours", COLOR_BLUE)
    _add_metric_card(slide, 3.65, 1.55, 2.55, "Total Cost", f"${metrics['demo_cost_m']:.2f}M", COLOR_GOLD)
    _add_metric_card(slide, 6.45, 1.55, 2.55, "Load Shedding", f"{metrics['demo_lol']:.1f}", COLOR_TEAL)
    _add_metric_card(slide, 9.25, 1.55, 2.55, "Dashboard Status", "Default case", COLOR_CORAL)

    chart_data = CategoryChartData()
    chart_data.categories = [zone for zone, _ in metrics["zone_costs"]]
    chart_data.add_series("Operating cost ($M)", [value for _, value in metrics["zone_costs"]])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.95),
        Inches(3.0),
        Inches(7.0),
        Inches(3.0),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Operating cost by zone"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = COLOR_NAVY
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = ZONE_COLORS[idx % len(ZONE_COLORS)]

    note = _add_card(slide, 8.35, 3.0, 3.7, 3.0, fill=RGBColor(236, 242, 250))
    note.text_frame.clear()
    lines = [
        "Why this matters",
        "The solved 2-day nodal case is stable enough to demo the Germany dashboard while the longer nodal horizon is still being staged.",
    ]
    for idx, line in enumerate(lines):
        p = note.text_frame.paragraphs[0] if idx == 0 else note.text_frame.add_paragraph()
        p.text = line
        _style_runs(p, 15 if idx == 0 else 18, COLOR_NAVY, bold=(idx == 0))
        if idx == 0:
            p.runs[0].font.color.rgb = COLOR_BLUE


def _add_next_steps_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Next Steps", "What to improve next in the Germany case build")

    cards = [
        ("Data QA", "Deepen plant validation and benchmark chronology totals against SMARD.", COLOR_BLUE),
        ("Nodal Scale-Up", "Extend from short-window nodal success toward longer production horizons.", COLOR_CORAL),
        ("Geometry", "Swap the reconstructed TSO overlay for official GIS if a public layer appears.", COLOR_TEAL),
        ("Storytelling", "Compare nodal and zonal outputs directly in the dashboard and in papers.", COLOR_PURPLE),
    ]
    positions = [(0.9, 1.7), (6.8, 1.7), (0.9, 4.05), (6.8, 4.05)]
    for (title, body, color), (left, top) in zip(cards, positions):
        card = _add_card(slide, left, top, 5.55, 1.8)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.18), Inches(1.8))
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.color.rgb = color
        title_box = _textbox(slide, left + 0.3, top + 0.22, 2.8, 0.25, title)
        _style_runs(title_box.text_frame.paragraphs[0], 17, color, bold=True)
        body_box = _textbox(slide, left + 0.3, top + 0.58, 4.8, 0.85, body)
        _style_runs(body_box.text_frame.paragraphs[0], 18, COLOR_NAVY)


def _add_repo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Key Files In The Repo", "Useful paths for updating the Germany build and the dashboard story")

    rows = [
        ("Build workspace", GERMANY_TOOLS.as_posix()),
        ("Dashboard geometry", GERMANY_DASHBOARD_GEOJSON.as_posix()),
        ("Nodal master case", (ROOT / "ModelCases" / "GERMANY_PCM_nodal_case").as_posix()),
        ("Zonal derivative case", (ROOT / "ModelCases" / "GERMANY_PCM_zonal4_case").as_posix()),
    ]
    top = 1.7
    for label, value in rows:
        card = _add_card(slide, 0.85, top, 11.65, 0.82)
        label_box = _textbox(slide, 1.08, top + 0.18, 2.2, 0.22, label)
        _style_runs(label_box.text_frame.paragraphs[0], 15, COLOR_BLUE, bold=True)
        value_box = _textbox(slide, 3.2, top + 0.16, 8.9, 0.26, value)
        _style_runs(value_box.text_frame.paragraphs[0], 14, COLOR_NAVY)
        top += 1.02


def _markdown_lines(metrics: dict[str, object]) -> list[str]:
    slides = [
        ("Germany PCM Case Build", ["From raw public data to consistent nodal and zonal HOPE cases."]),
        ("Executive Summary", [
            "Goal: create one canonical Germany PCM dataset that supports both nodal and zonal analysis.",
            "Method: freeze every major mapping once, then derive the zonal case from the nodal case.",
            "Outcome: solved 2-day nodal demo case plus a consistent 4-zone zonal comparison case.",
        ]),
        ("Source Stack", [
            "Network backbone: OSM Europe transmission dataset plus PyPSA-Eur workflow reference.",
            "Fleet: powerplantmatching with BNetzA validation layers.",
            "Chronology: SMARD national load and generation plus four TSO load helper files.",
            "Map geometry: Germany state-boundary GeoJSON reconstructed into a dashboard TSO layer.",
        ]),
        ("Integration Architecture", [
            "Raw sources feed cleaned staging tables, frozen maps, the nodal master case, and the zonal derivative case.",
            "The zonal case is produced by aggregation from the nodal master case.",
        ]),
        ("Build Steps", [
            "1. Network backbone",
            "2. Bus-zone mapping",
            "3. Generator mapping",
            "4. Chronology",
            "5. Case assembly",
            "6. Dashboard geometry",
        ]),
        ("Case Snapshot", [
            f"Nodal master case: {metrics['nodal_buses']} buses, {metrics['nodal_lines']} lines, {metrics['nodal_gens']} generators.",
            f"Zonal derivative case: {metrics['zonal_zones']} zones, {metrics['zonal_lines']} interfaces, {metrics['zonal_gens']} generators.",
        ]),
        ("Assumptions And Caveats", [
            "The 4-zone Germany setup is a research zoning, not the real DE-LU bidding zone.",
            "The dashboard TSO overlay is reconstructed, not an official public shapefile.",
            "The nodal case is the canonical truth source and the zonal case is derived from it.",
        ]),
        ("Debugging Journey", [
            "Initial nodal failures were caused by topology and reactance-scaling issues.",
            "Fixes: transformer connectivity, largest connected component, and HOPE-scale reactance normalization.",
            "Result: stable 1-day and 2-day full-nodal debug runs.",
        ]),
        ("Solved Demo Outcome", [
            f"2-day January nodal demo total cost: ${metrics['demo_cost_m']:.2f}M.",
            f"Load shedding: {metrics['demo_lol']:.1f}.",
            "The 2-day nodal case is now the dashboard default.",
        ]),
        ("Next Steps", [
            "Deepen validation and benchmarking.",
            "Extend nodal solve horizon.",
            "Upgrade geometry if official GIS becomes available.",
            "Strengthen nodal-vs-zonal storytelling in the dashboard.",
        ]),
        ("Key Files In The Repo", [
            GERMANY_TOOLS.as_posix(),
            GERMANY_DASHBOARD_GEOJSON.as_posix(),
            (ROOT / "ModelCases" / "GERMANY_PCM_nodal_case").as_posix(),
            (ROOT / "ModelCases" / "GERMANY_PCM_zonal4_case").as_posix(),
        ]),
    ]
    lines = ["# Germany PCM Case Build Deck", ""]
    for idx, (title, bullets) in enumerate(slides, start=1):
        lines.append(f"## Slide {idx}. {title}")
        for bullet in bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    return lines


def build_presentation() -> None:
    metrics = _case_metrics()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, metrics)
    _add_summary_slide(prs)
    _add_source_slide(prs)
    _add_architecture_slide(prs)
    _add_steps_slide(prs, metrics)
    _add_snapshot_slide(prs, metrics)
    _add_assumptions_slide(prs)
    _add_debug_slide(prs)
    _add_results_slide(prs, metrics)
    _add_next_steps_slide(prs)
    _add_repo_slide(prs)

    prs.save(OUTPUT_PPTX)
    OUTPUT_MD.write_text("\n".join(_markdown_lines(metrics)), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
