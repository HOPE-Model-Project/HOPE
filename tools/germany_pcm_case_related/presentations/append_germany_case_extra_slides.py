from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
PRESENTATION_DIR = Path(__file__).resolve().parent
PPTX_PATH = PRESENTATION_DIR / "GERMANY_PCM_CASE_BUILD_DECK.pptx"
MD_PATH = PRESENTATION_DIR / "GERMANY_PCM_CASE_BUILD_DECK.md"

RAW_ROOT = ROOT / "tools" / "germany_pcm_case_related" / "raw_sources"
REF_ROOT = ROOT / "tools" / "germany_pcm_case_related" / "references"

COLOR_BG = RGBColor(246, 248, 251)
COLOR_PANEL = RGBColor(255, 255, 255)
COLOR_NAVY = RGBColor(18, 42, 76)
COLOR_BLUE = RGBColor(47, 93, 169)
COLOR_TEAL = RGBColor(54, 143, 132)
COLOR_GOLD = RGBColor(219, 165, 32)
COLOR_CORAL = RGBColor(225, 107, 74)
COLOR_PURPLE = RGBColor(145, 96, 192)
COLOR_SLATE = RGBColor(91, 105, 125)
COLOR_LINE = RGBColor(214, 223, 233)
COLOR_WHITE = RGBColor(255, 255, 255)


EXTRA_SLIDE_TITLES = [
    "Raw Input Snapshot",
    "Processing To The Current Case",
    "Open-Source Case Comparison",
    "Comparison: Pros And Trade-Offs",
]


def _safe_csv_rows(path: Path) -> int:
    if not path.exists():
        return 0
    try:
        return len(pd.read_csv(path, sep=None, engine="python"))
    except Exception:
        with path.open("r", encoding="utf-8", errors="ignore") as handle:
            return max(sum(1 for _ in handle) - 1, 0)


def _metrics() -> dict[str, int]:
    return {
        "raw_buses": _safe_csv_rows(RAW_ROOT / "osm_europe_grid" / "buses.csv"),
        "raw_lines": _safe_csv_rows(RAW_ROOT / "osm_europe_grid" / "lines.csv"),
        "raw_transformers": _safe_csv_rows(RAW_ROOT / "osm_europe_grid" / "transformers.csv"),
        "raw_links": _safe_csv_rows(RAW_ROOT / "osm_europe_grid" / "links.csv"),
        "raw_fleet": _safe_csv_rows(RAW_ROOT / "powerplantmatching" / "germany_powerplantmatching.csv"),
        "raw_generation_hours": _safe_csv_rows(RAW_ROOT / "smard_2025" / "germany_actual_generation_hourly.csv"),
        "raw_tso_hours": _safe_csv_rows(RAW_ROOT / "smard_2025" / "load_50Hertz_hourly.csv"),
        "clean_buses": len(pd.read_csv(REF_ROOT / "germany_network_buses_clean.csv")),
        "clean_lines": len(pd.read_csv(REF_ROOT / "germany_network_lines_clean.csv")),
        "clean_transformers": len(pd.read_csv(REF_ROOT / "germany_network_transformers_clean.csv")),
        "clean_fleet": len(pd.read_csv(REF_ROOT / "germany_generator_fleet_clean.csv")),
        "clean_bus_zone": len(pd.read_csv(REF_ROOT / "germany_bus_zone_map.csv")),
        "clean_hours": len(pd.read_csv(REF_ROOT / "germany_hourly_chronology_clean.csv")),
    }


def _existing_titles(prs: Presentation) -> set[str]:
    titles: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text_frame.text.strip()
            if text:
                titles.add(text.split("\n")[0].strip())
    return titles


def _set_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG


def _textbox(slide, left: float, top: float, width: float, height: float, text: str = ""):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if text:
        box.text_frame.text = text
    return box


def _style_runs(paragraph, size: int, color: RGBColor, bold: bool = False) -> None:
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = _textbox(slide, 0.6, 0.38, 11.8, 0.55, title)
    _style_runs(title_box.text_frame.paragraphs[0], 26, COLOR_NAVY, bold=True)
    if subtitle:
        sub_box = _textbox(slide, 0.62, 0.96, 11.5, 0.32, subtitle)
        _style_runs(sub_box.text_frame.paragraphs[0], 13, COLOR_SLATE)


def _add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = COLOR_PANEL, line: RGBColor = COLOR_LINE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def _add_metric_card(slide, left: float, top: float, width: float, label: str, value: str, accent: RGBColor) -> None:
    card = _add_card(slide, left, top, width, 1.0, fill=COLOR_WHITE)
    card.line.color.rgb = accent
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.14), Inches(1.0))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent
    label_box = _textbox(slide, left + 0.24, top + 0.14, width - 0.34, 0.25, label)
    _style_runs(label_box.text_frame.paragraphs[0], 11, COLOR_SLATE, bold=True)
    value_box = _textbox(slide, left + 0.24, top + 0.42, width - 0.34, 0.34, value)
    _style_runs(value_box.text_frame.paragraphs[0], 23, COLOR_NAVY, bold=True)


def _add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], size: int = 18) -> None:
    box = _textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.space_after = Pt(8)
        _style_runs(p, size, COLOR_NAVY)


def _append_raw_input_slide(prs: Presentation, metrics: dict[str, int]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Raw Input Snapshot", "What went into the Germany build before cleaning and integration")

    raw_cards = [
        ("Raw buses", f"{metrics['raw_buses']:,}", COLOR_BLUE),
        ("Raw lines", f"{metrics['raw_lines']:,}", COLOR_CORAL),
        ("Raw transformers", f"{metrics['raw_transformers']:,}", COLOR_TEAL),
        ("Raw fleet rows", f"{metrics['raw_fleet']:,}", COLOR_GOLD),
    ]
    for idx, (label, value, color) in enumerate(raw_cards):
        _add_metric_card(slide, 0.82 + idx * 3.05, 1.55, 2.72, label, value, color)

    clean = _add_card(slide, 0.9, 3.0, 5.9, 2.7)
    clean.text_frame.clear()
    for idx, line in enumerate(
        [
            "Cleaned staging outputs",
            f"{metrics['clean_buses']:,} cleaned network buses and {metrics['clean_lines']:,} cleaned lines",
            f"{metrics['clean_transformers']:,} cleaned transformers and {metrics['clean_bus_zone']:,} bus-zone assignments",
            f"{metrics['clean_fleet']:,} cleaned generator rows and {metrics['clean_hours']:,} canonical hourly records",
        ]
    ):
        p = clean.text_frame.paragraphs[0] if idx == 0 else clean.text_frame.add_paragraph()
        p.text = line
        _style_runs(p, 16 if idx else 17, COLOR_NAVY, bold=(idx == 0))
        if idx == 0:
            p.runs[0].font.color.rgb = COLOR_BLUE

    notes = [
        f"SMARD generation file: {metrics['raw_generation_hours']:,} raw rows before trimming to the canonical 8,760-hour basis.",
        f"Each SMARD TSO load helper file: {metrics['raw_tso_hours']:,} raw rows before chronology normalization.",
        "The raw-to-clean reduction is intentional: the HOPE case keeps only Germany-relevant, schema-consistent, and solve-ready records.",
    ]
    _add_bullets(slide, 7.15, 3.0, 5.0, 2.8, notes, size=17)


def _append_processing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Processing To The Current Case", "How the raw files become the current Germany nodal and zonal HOPE cases")

    steps = [
        ("1", "Ingest raw sources", "OSM network tables, powerplantmatching fleet, SMARD chronology, and reference geometry."),
        ("2", "Normalize schemas", "Rename columns, coerce coordinates, standardize technology and status fields, and filter to Germany."),
        ("3", "Repair the network", "Keep transformer connectivity, remove disconnected leftovers, and rescale branch reactance for HOPE."),
        ("4", "Freeze spatial maps", "Create one Bus_id -> Zone_id table and one generator -> bus mapping table."),
        ("5", "Build chronology", "Normalize Germany hourly load and generation, then create zonal helper load shares."),
        ("6", "Assemble cases", "Build the nodal master case first, then derive the 4-zone case and the dashboard geometry from it."),
    ]
    top = 1.55
    for num, title, body in steps:
        bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.95), Inches(top), Inches(0.5), Inches(0.5))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = COLOR_NAVY
        bubble.line.color.rgb = COLOR_NAVY
        bubble.text_frame.text = num
        bubble.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _style_runs(bubble.text_frame.paragraphs[0], 17, COLOR_WHITE, bold=True)
        t = _textbox(slide, 1.65, top - 0.02, 3.2, 0.24, title)
        _style_runs(t.text_frame.paragraphs[0], 18, COLOR_BLUE, bold=True)
        b = _textbox(slide, 1.66, top + 0.22, 10.1, 0.36, body)
        _style_runs(b.text_frame.paragraphs[0], 16, COLOR_NAVY)
        top += 0.86


def _append_comparison_overview_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Open-Source Case Comparison", "Where this Germany HOPE case sits relative to other open-source electricity-model case frameworks")

    table = slide.shapes.add_table(5, 4, Inches(0.75), Inches(1.55), Inches(11.9), Inches(4.8)).table
    headers = ["Case / framework", "Primary focus", "Geographic style", "Why it matters here"]
    rows = [
        ["This HOPE Germany case", "Consistent nodal-vs-zonal PCM comparison", "Germany nodal master plus derived 4-zone case", "Built specifically to compare dispatch, congestion, and price outcomes under matched assumptions"],
        ["PyPSA-Eur", "Reproducible Europe-wide workflow and scenario engine", "Europe-wide network with flexible clustering", "Best backbone reference for open data retrieval and preprocessing structure"],
        ["eTraGo / open_eGo", "High-resolution Germany-centered planning and sector coupling", "Germany detailed, foreign system aggregated", "Strong Germany data model and nodal optimization context"],
        ["POMATO DE example", "Market design, FBMC, redispatch, and zonal analysis", "Germany / European market-study orientation", "Strong reference for zonal market-coupling and redispatch analysis"],
    ]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        for p in cell.text_frame.paragraphs:
            _style_runs(p, 12, COLOR_WHITE, bold=True)
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 else RGBColor(240, 244, 249)
            for p in cell.text_frame.paragraphs:
                _style_runs(p, 11, COLOR_NAVY)
                p.alignment = PP_ALIGN.LEFT


def _append_comparison_tradeoff_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Comparison: Pros And Trade-Offs", "Why we still built a custom HOPE Germany case instead of using one of the open models directly")

    cards = [
        (
            "HOPE Germany case",
            ["Pros", "Direct nodal-vs-zonal comparability", "Same chronology and asset assumptions in both cases"],
            ["Trade-offs", "Short nodal solve horizon is solved and demo-ready, but longer horizons still need staging"],
            COLOR_BLUE,
        ),
        (
            "PyPSA-Eur",
            ["Pros", "Excellent open workflow, clustering, and data-pipeline reference"],
            ["Trade-offs", "Not a drop-in HOPE PCM case and not tailored to Germany 4-TSO nodal-vs-zonal comparison"],
            COLOR_CORAL,
        ),
        (
            "eTraGo / open_eGo",
            ["Pros", "Rich Germany-specific data model and strong nodal / planning context"],
            ["Trade-offs", "Heavier stack, database dependence, and broader scope than the focused HOPE comparison case"],
            COLOR_TEAL,
        ),
        (
            "POMATO DE example",
            ["Pros", "Strong market-coupling, FBMC, and redispatch framing"],
            ["Trade-offs", "Different modeling stack and case objective than our HOPE PCM comparison build"],
            COLOR_PURPLE,
        ),
    ]
    positions = [(0.85, 1.7), (6.8, 1.7), (0.85, 4.05), (6.8, 4.05)]
    for (title, pros, cons, color), (left, top) in zip(cards, positions):
        card = _add_card(slide, left, top, 5.55, 1.8)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(5.55), Inches(0.16))
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.color.rgb = color
        t = _textbox(slide, left + 0.22, top + 0.22, 3.6, 0.24, title)
        _style_runs(t.text_frame.paragraphs[0], 16, color, bold=True)
        pbox = _textbox(slide, left + 0.24, top + 0.58, 4.95, 0.9)
        tf = pbox.text_frame
        tf.clear()
        for idx, line in enumerate(pros + cons):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            _style_runs(p, 13 if idx in {0, len(pros)} else 14, COLOR_NAVY, bold=(idx in {0, len(pros)}))
            if idx in {0, len(pros)}:
                p.runs[0].font.color.rgb = color


def _append_md_sections() -> None:
    text = MD_PATH.read_text(encoding="utf-8") if MD_PATH.exists() else "# Germany PCM Case Build Deck\n\n"
    for title in EXTRA_SLIDE_TITLES:
        if title in text:
            return
    extra = """
## Extra Slide. Raw Input Snapshot
- Raw network files: 12,936 buses, 16,050 line rows, 1,949 transformers, 35 links.
- Raw fleet file: 165,064 Germany plant rows from powerplantmatching.
- Raw SMARD inputs: 8,766 rows in the national generation file and 8,766 rows in each TSO load helper file before normalization.
- Clean staging outputs: 791 buses, 1,019 lines, 155 transformers, 141,420 cleaned generator rows, and 8,760 canonical chronology rows.

## Extra Slide. Processing To The Current Case
- Ingest raw network, fleet, chronology, and reference geometry files.
- Normalize schemas and filter to Germany-relevant records.
- Repair topology and electrical scaling for HOPE compatibility.
- Freeze bus-zone and generator-bus mapping tables.
- Normalize chronology and zonal helper shares.
- Build the nodal master case first and derive the zonal case from it.

## Extra Slide. Open-Source Case Comparison
- HOPE Germany case: focused on matched nodal-vs-zonal PCM comparison.
- PyPSA-Eur: strong Europe-wide workflow and clustering reference.
- eTraGo/open_eGo: strong Germany-specific nodal and planning context.
- POMATO DE example: strong market-coupling and redispatch framing.

## Extra Slide. Comparison: Pros And Trade-Offs
- HOPE Germany case: best for direct comparison under consistent assumptions, but still being scaled to longer nodal horizons.
- PyPSA-Eur: best open workflow backbone, but not a drop-in HOPE PCM case.
- eTraGo/open_eGo: rich German system detail, but heavier and broader than the focused HOPE comparison build.
- POMATO DE example: very strong for zonal market design, but a different modeling stack and study objective.
"""
    MD_PATH.write_text(text.rstrip() + "\n\n" + extra.strip() + "\n", encoding="utf-8")


def append_slides() -> None:
    prs = Presentation(PPTX_PATH)
    titles = _existing_titles(prs)
    if "Raw Input Snapshot" not in titles:
        _append_raw_input_slide(prs, _metrics())
    if "Processing To The Current Case" not in titles:
        _append_processing_slide(prs)
    if "Open-Source Case Comparison" not in titles:
        _append_comparison_overview_slide(prs)
    if "Comparison: Pros And Trade-Offs" not in titles:
        _append_comparison_tradeoff_slide(prs)
    prs.save(PPTX_PATH)
    _append_md_sections()


if __name__ == "__main__":
    append_slides()
